from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement

def add_shape_with_text(slide, left, top, width, height, text, font_size=18, bold=False, color=RGBColor(0, 0, 0), align=PP_ALIGN.LEFT):
    """Helper function to add text box"""
    shape = slide.shapes.add_textbox(left, top, width, height)
    text_frame = shape.text_frame
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return shape

def load_template_and_create():
    """Load template and create AEGIS presentation"""
    
    # Load template
    template_path = "Smart-Horizon-2026-48Hour-Intrnl Hackathon Level-1 Template.pptx"
    prs = Presentation(template_path)
    
    # Slide 1: Title Slide
    slide1 = prs.slides[0]
    for shape in slide1.shapes:
        if hasattr(shape, "text") and "Enter your proposed solution title here" in shape.text:
            shape.text = "AEGIS v1\nAutonomous Swarm Control for Disaster Recovery"
            for paragraph in shape.text_frame.paragraphs:
                paragraph.font.size = Pt(40)
                paragraph.font.bold = True
                paragraph.font.color.rgb = RGBColor(0, 102, 204)
    
    # Slide 2: Problem Understanding
    slide2 = prs.slides[1]
    for shape in slide2.shapes:
        if hasattr(shape, "text") and "Problem description" in shape.text:
            content = """• Disaster scenarios (earthquakes, tsunamis, wildfires, floods) require rapid response
• Manual search and rescue operations are time-consuming and dangerous
• Need for autonomous, decentralized drone swarms for survivor detection
• Current solutions lack real-time coordination and adaptive path planning
• Hardware constraints (battery, communication range, sensor limits) not modeled"""
            shape.text = content
            for paragraph in shape.text_frame.paragraphs:
                paragraph.font.size = Pt(14)
                paragraph.space_before = Pt(6)
    
    # Slide 3: Proposed Solution & Innovation
    slide3 = prs.slides[2]
    for shape in slide3.shapes:
        if hasattr(shape, "text") and "Overview of proposed solution" in shape.text:
            content = """• AEGIS: Real-time mission control dashboard for autonomous drone swarms
• Multi-scenario simulation (earthquake, tsunami, wildfire, flood)
• Key Features:
  - Decentralized control with dynamic zone allocation
  - Real-time 3D visualization (Three.js + WebGL)
  - Thermal imaging + camera feed integration
  - A* pathfinding with survivor detection
  - Battery management & charging station logic
• Innovation: Combines physics simulation with low-latency tactical interface"""
            shape.text = content
            for paragraph in shape.text_frame.paragraphs:
                paragraph.font.size = Pt(13)
                paragraph.space_before = Pt(4)
    
    # Slide 4: Technical Architecture
    slide4 = prs.slides[3]
    for shape in slide4.shapes:
        if hasattr(shape, "text") and "System architecture" in shape.text:
            content = """Technology Stack:
• FRONTEND: Vite + React + Three.js (Fiber/Drei) + Zustand (state management)
• BACKEND: Python 3.x + FastAPI (async WebSocket)
• SIMULATION: NumPy-based physics engine + A* pathfinding
• VISUALIZATION: Real-time 3D terrain, fire/smoke, drone trails

Architecture:
• Decentralized control: Each drone operates autonomously with assigned zones
• WebSocket hub: Real-time state sync (state updates @ 20Hz)
• Physics loop: 50ms tick cycle with battery/movement calculations
• Event logging: Comprehensive mission event tracking

Technologies: FastAPI, WebSockets, NumPy, Three.js, React, Tailwind CSS"""
            shape.text = content
            for paragraph in shape.text_frame.paragraphs:
                paragraph.font.size = Pt(12)
                paragraph.space_before = Pt(3)
    
    # Slide 5: Implementation / Prototype
    slide5 = prs.slides[4]
    for shape in slide5.shapes:
        if hasattr(shape, "text") and "Prototype screenshots" in shape.text:
            content = """Working Features:
• Multi-drone tracking: 5 autonomous drones with real-time telemetry
• 3D disaster scene: Procedural urban terrain with fire/smoke simulation
• 2D tactical map: "Treasure map" styled long-range surveyor plotting
• Thermal imaging: Live thermal camera feed with survivor heat signatures
• Survivor detection: Probabilistic detection model with confidence scoring
• Event log: Real-time mission events (warnings, detections, status updates)

Key Functionalities:
• 4 disaster scenarios fully operational
• Live battery management & charging station logic
• Orbit-based drone physics with micro-turbulence
• Hazard zone avoidance (gas leaks, flooding, fire)
• Web-based dashboard accessible from any device"""
            shape.text = content
            for paragraph in shape.text_frame.paragraphs:
                paragraph.font.size = Pt(12)
                paragraph.space_before = Pt(3)
    
    # Slide 6: Feasibility & Impact
    slide6 = prs.slides[5]
    for shape in slide6.shapes:
        if hasattr(shape, "text") and "Real-world applicability" in shape.text:
            content = """Real-world Applicability:
✓ Deployable on DJI Matrice 300 RTK (baseline hardware)
✓ Scalable to 10+ drones with proper load balancing
✓ Applicable to earthquake, wildfire, flood, tsunami scenarios

Social Impact:
• Faster search & rescue operations (minutes vs. hours)
• Reduced risk to human rescue teams
• Improved survivor detection rates in disaster zones
• Reusable simulation for training & planning

Future Enhancements (Phase 1-4):
• Implement RRT* 3D path planning with dynamic collision avoidance
• Add realistic battery model (DJI Matrice 300 RTK specs)
• Integrate ORB-SLAM2 for visual SLAM
• Add YOLO object detection for autonomous threat identification
• Hardware-in-the-loop simulation for real drone integration"""
            shape.text = content
            for paragraph in shape.text_frame.paragraphs:
                paragraph.font.size = Pt(12)
                paragraph.space_before = Pt(3)
    
    # Slide 7: Conclusion
    slide7 = prs.slides[6]
    for shape in slide7.shapes:
        if hasattr(shape, "text") and "Summary of the solution" in shape.text:
            content = """Summary:
AEGIS v1 demonstrates a solid decentralized swarm architecture with real-time 3D visualization and multi-scenario disaster recovery simulation.

Key Outcomes:
✓ Fully operational mission control dashboard
✓ Physics-based drone simulation (5 drones, 4 scenarios)
✓ Real-time thermal imaging & survivor detection
✓ Comprehensive event logging & telemetry

Roadmap:
Phase 1: Add research references + hardware constraints (1-2 weeks)
Phase 2: Implement dynamic task allocation + RRT* planning (2-4 weeks)
Phase 3: SLAM + YOLO integration + fault tolerance (1-2 months)
Phase 4: Hardware trials + field deployment (ongoing)

Team Status: Ready for Phase 1 implementation"""
            shape.text = content
            for paragraph in shape.text_frame.paragraphs:
                paragraph.font.size = Pt(12)
                paragraph.space_before = Pt(3)
    
    # Slide 8: Thank You (Keep as is or update)
    slide8 = prs.slides[7]
    
    # Save presentation
    output_path = "AEGIS_Presentation_SmartHorizon2026.pptx"
    prs.save(output_path)
    print(f"✓ Presentation created: {output_path}")
    return output_path

if __name__ == "__main__":
    try:
        output = load_template_and_create()
        print(f"✓ Successfully generated: {output}")
        print(f"✓ 8 slides with AEGIS project content")
        print(f"✓ Based on Smart Horizon 2026 template structure")
    except Exception as e:
        print(f"Error: {e}")
        import traceback
        traceback.print_exc()
