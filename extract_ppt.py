from pptx import Presentation
import os

# Load the template
template_path = "Smart-Horizon-2026-48Hour-Intrnl Hackathon Level-1 Template.pptx"

try:
    prs = Presentation(template_path)
    
    print("=" * 80)
    print(f"TEMPLATE ANALYSIS: {os.path.basename(template_path)}")
    print(f"Total Slides: {len(prs.slides)}")
    print("=" * 80)
    
    for slide_num, slide in enumerate(prs.slides, 1):
        print(f"\n{'='*80}")
        print(f"SLIDE {slide_num}")
        print(f"{'='*80}")
        print(f"Layout: {slide.slide_layout.name}")
        print(f"Shapes: {len(slide.shapes)}")
        
        for shape_idx, shape in enumerate(slide.shapes):
            if hasattr(shape, "text"):
                text_content = shape.text.strip()
                if text_content:
                    print(f"\n  Shape {shape_idx} ({shape.shape_type}):")
                    print(f"    Text: {text_content[:300]}")
            
            if hasattr(shape, "text_frame"):
                for para_idx, paragraph in enumerate(shape.text_frame.paragraphs):
                    if paragraph.text.strip():
                        print(f"    Paragraph {para_idx}: {paragraph.text[:150]}")

except Exception as e:
    print(f"Error: {e}")
    import traceback
    traceback.print_exc()
